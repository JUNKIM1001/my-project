import io
import os
import re

import fitz  # PyMuPDF
import requests
from bs4 import BeautifulSoup
from google import genai
from google.genai import types

from pypdf import PdfReader
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

HEADERS = {
    "User-Agent": "Mozilla/5.0 (compatible; LeadFinderBot/0.1)"
}

# テキスト抽出がこの文字数未満なら画像主体の資料とみなし、画像認識にフォールバック
MIN_TEXT_LEN = 80
MAX_IMAGES = 25  # マルチモーダルに送るスライド画像の上限

MODEL = "gemini-flash-latest"

SUMMARY_PROMPT = """あなたはB2B営業の企画を支援するアナリストです。
以下はある会社の営業資料・サービス資料から抽出したテキストです。
この会社の「自社商材情報」を、営業ターゲット（事業会社）の探索に使えるように日本語で簡潔に構造化してください。

次の見出しで、わかる範囲だけを箇条書きでまとめてください（資料に無い項目は省略）:
- 会社名 / サービス名
- 商材の内容（何を提供しているか、1〜2文）
- 解決できる顧客の課題（この商材が刺さる状況）
- 想定顧客・ターゲット市場（業界・企業規模・部署など）
- 導入実績・強み（事例、差別化ポイント）
- 価格帯・提供形態（わかる場合のみ）

説明文や前置きは不要で、見出し付きの箇条書きのみを出力してください。
"""


def extract_text_from_pdf(data: bytes) -> str:
    reader = PdfReader(io.BytesIO(data))
    parts = []
    for page in reader.pages:
        try:
            parts.append(page.extract_text() or "")
        except Exception:
            continue
    return "\n".join(parts)


def extract_text_from_pptx(data: bytes) -> str:
    prs = Presentation(io.BytesIO(data))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if text.strip():
                        parts.append(text)
    return "\n".join(parts)


def extract_text(filename: str, data: bytes) -> str:
    lower = filename.lower()
    if lower.endswith(".pdf"):
        return extract_text_from_pdf(data)
    if lower.endswith(".pptx"):
        return extract_text_from_pptx(data)
    raise ValueError("対応形式は PDF または PPTX です。")


def _walk_shapes(shapes):
    """グループ図形も再帰的に展開して全図形を返す。"""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk_shapes(shape.shapes)
        else:
            yield shape


def extract_images_from_pptx(data: bytes) -> list:
    """PPTXに埋め込まれた画像（スライド画像）のバイト列リストを返す。"""
    prs = Presentation(io.BytesIO(data))
    images = []
    for slide in prs.slides:
        for shape in _walk_shapes(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    images.append(shape.image.blob)
                except Exception:
                    continue
        if len(images) >= MAX_IMAGES:
            break
    return images[:MAX_IMAGES]


def render_pdf_to_images(data: bytes) -> list:
    """PDFの各ページをPNG画像にレンダリングしてバイト列リストを返す。"""
    doc = fitz.open(stream=data, filetype="pdf")
    images = []
    for page in doc:
        pix = page.get_pixmap(dpi=120)
        images.append(pix.tobytes("png"))
        if len(images) >= MAX_IMAGES:
            break
    return images


def extract_images(filename: str, data: bytes) -> list:
    lower = filename.lower()
    if lower.endswith(".pdf"):
        return render_pdf_to_images(data)
    if lower.endswith(".pptx"):
        return extract_images_from_pptx(data)
    return []


def extract_text_from_html(html: str) -> str:
    soup = BeautifulSoup(html, "lxml")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    text = soup.get_text(separator=" ")
    return re.sub(r"\s+", " ", text).strip()


def _google_export_url(url: str) -> str:
    """Google Slides/Docs の共有URLを、本文を取得できるエクスポートURLに変換する。"""
    m = re.search(r"docs\.google\.com/(presentation|document)/d/([a-zA-Z0-9_-]+)", url)
    if not m:
        return url
    kind, doc_id = m.group(1), m.group(2)
    if kind == "presentation":
        # スライドはPDFとしてエクスポートするとテキストを抽出できる
        return f"https://docs.google.com/presentation/d/{doc_id}/export/pdf"
    return f"https://docs.google.com/document/d/{doc_id}/export?format=txt"


def extract_text_from_url(url: str) -> str:
    """URL（Google Slides共有URL / PDF / PPTX / 一般Webページ）からテキストを抽出する。"""
    fetch_url = _google_export_url(url)
    try:
        resp = requests.get(fetch_url, headers=HEADERS, timeout=30)
        resp.raise_for_status()
    except Exception as e:
        raise ValueError(f"URLの取得に失敗しました（共有設定が「リンクを知っている全員」になっているか確認してください）: {e}")

    content_type = (resp.headers.get("Content-Type") or "").lower()
    lower_url = fetch_url.lower()

    if "application/pdf" in content_type or lower_url.endswith("/export/pdf") or lower_url.endswith(".pdf"):
        return extract_text_from_pdf(resp.content)
    if "presentation" in content_type or lower_url.endswith(".pptx"):
        return extract_text_from_pptx(resp.content)
    if "text/plain" in content_type or "format=txt" in lower_url:
        return resp.content.decode("utf-8", errors="ignore")
    # それ以外は一般的なWebページとして本文を抽出
    return extract_text_from_html(resp.text)


def summarize_company(raw_text: str) -> dict:
    api_key = os.environ.get("GEMINI_API_KEY")
    if not api_key:
        return {"error": "GEMINI_API_KEY が設定されていません。.env を確認してください。"}

    text = raw_text.strip()
    if not text:
        return {"error": "資料からテキストを抽出できませんでした（画像のみのPDF等の可能性があります）。"}

    client = genai.Client(api_key=api_key)
    try:
        response = client.models.generate_content(
            model=MODEL,
            contents=f"{SUMMARY_PROMPT}\n\n--- 営業資料テキスト ---\n{text[:20000]}",
        )
    except Exception as e:
        return {"error": f"要約の生成に失敗しました: {e}"}

    return {"summary": (response.text or "").strip()}


def summarize_company_from_images(images: list) -> dict:
    """スライド画像（画像主体の資料）をGeminiの画像認識で読み取り要約する。"""
    api_key = os.environ.get("GEMINI_API_KEY")
    if not api_key:
        return {"error": "GEMINI_API_KEY が設定されていません。.env を確認してください。"}
    if not images:
        return {"error": "資料から画像を取得できませんでした。"}

    client = genai.Client(api_key=api_key)
    parts = [types.Part.from_text(text=SUMMARY_PROMPT + "\n\n以下は営業資料の各スライド画像です。")]
    for img in images:
        parts.append(types.Part.from_bytes(data=img, mime_type="image/png"))

    try:
        response = client.models.generate_content(model=MODEL, contents=parts)
    except Exception as e:
        return {"error": f"画像からの要約生成に失敗しました: {e}"}

    return {"summary": (response.text or "").strip()}


def build_profile_from_file(filename: str, data: bytes) -> dict:
    """ファイル（PDF/PPTX）から自社事業情報を生成。テキストが乏しければ画像認識にフォールバック。"""
    try:
        raw_text = extract_text(filename, data)
    except ValueError as e:
        return {"error": str(e)}

    if len(raw_text.strip()) >= MIN_TEXT_LEN:
        result = summarize_company(raw_text)
        if "summary" in result:
            result["raw_text"] = raw_text
        return result

    # テキストが乏しい → 画像主体の資料とみなして画像認識
    images = extract_images(filename, data)
    result = summarize_company_from_images(images)
    if "summary" in result:
        result["raw_text"] = raw_text or "(画像認識により抽出)"
    return result


def build_profile_from_url(url: str) -> dict:
    """URLから自社事業情報を生成。Google Slides等の画像主体資料は画像認識にフォールバック。"""
    fetch_url = _google_export_url(url)
    try:
        resp = requests.get(fetch_url, headers=HEADERS, timeout=30)
        resp.raise_for_status()
    except Exception as e:
        return {"error": f"URLの取得に失敗しました（共有設定が「リンクを知っている全員」か確認してください）: {e}"}

    content_type = (resp.headers.get("Content-Type") or "").lower()
    lower_url = fetch_url.lower()
    is_pdf = "application/pdf" in content_type or lower_url.endswith("/export/pdf") or lower_url.endswith(".pdf")
    is_pptx = "presentation" in content_type or lower_url.endswith(".pptx")

    if is_pdf:
        raw_text = extract_text_from_pdf(resp.content)
        images_fn = lambda: render_pdf_to_images(resp.content)
    elif is_pptx:
        raw_text = extract_text_from_pptx(resp.content)
        images_fn = lambda: extract_images_from_pptx(resp.content)
    elif "text/plain" in content_type or "format=txt" in lower_url:
        raw_text = resp.content.decode("utf-8", errors="ignore")
        images_fn = lambda: []
    else:
        raw_text = extract_text_from_html(resp.text)
        images_fn = lambda: []

    if len(raw_text.strip()) >= MIN_TEXT_LEN:
        result = summarize_company(raw_text)
        if "summary" in result:
            result["raw_text"] = raw_text
        return result

    images = images_fn()
    if not images:
        return {"error": "資料からテキスト・画像を抽出できませんでした。共有設定や資料の内容をご確認ください。"}
    result = summarize_company_from_images(images)
    if "summary" in result:
        result["raw_text"] = raw_text or "(画像認識により抽出)"
    return result
